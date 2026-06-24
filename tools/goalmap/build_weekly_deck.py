#!/usr/bin/env python3
"""週次MTG用スライド（PPTX）を生成する。

members/*.json を全員分読み込み、
  1) 表紙＋全員ダッシュボード（達成率・現在ステージ一覧）
  2) 各メンバー1枚（ゴール図解＋今週のフォーカス／ネクスト）
を1つの .pptx にまとめる。図解は generate_goalmap.py と同一仕様。

使い方:
    python tools/goalmap/build_weekly_deck.py                 # members/*.json 全員
    python tools/goalmap/build_weekly_deck.py 岡野 宮腰        # 指定メンバーのみ
    python tools/goalmap/build_weekly_deck.py -o /path/deck.pptx

要件: pip install python-pptx cairosvg / CJKフォント（fonts-noto-cjk）
"""
from __future__ import annotations

import datetime as dt
import sys
import tempfile
from pathlib import Path

import cairosvg  # type: ignore
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parent))
from generate_goalmap import FONT_RASTER, build_svg, load_member_data  # noqa: E402

JP = "Hiragino Sans"
INK = RGBColor(0x1F, 0x29, 0x33)
SUB = RGBColor(0x67, 0x70, 0x7A)
DONE = RGBColor(0x1D, 0x9E, 0x75)
NOW = RGBColor(0xEF, 0x9F, 0x27)
GOAL = RGBColor(0x63, 0x57, 0xCC)
FUT_BG = RGBColor(0xF2, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HERE = Path(__file__).parent

STAGE_NAMES = ["①型を知る", "②練習", "③実践", "④振り返り", "⑤自走"]


def _font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = JP
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", JP)


def _text(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = txt
        _font(r, size, color, bold)
    return tb


def _rate(d) -> int:
    tasks = [t for p in d["phases"] for t in p.get("tasks", [])]
    done = sum(1 for t in tasks if t.get("done"))
    return round(done / len(tasks) * 100) if tasks else 0


def _png(d) -> bytes:
    svg = build_svg(d, font=FONT_RASTER)
    return cairosvg.svg2png(bytestring=svg.encode("utf-8"), scale=2)


def _png_size(d) -> tuple[int, int]:
    """図解の論理サイズ(px) = 680 × H。"""
    import re
    svg = build_svg(d)
    m = re.search(r'viewBox="0 0 (\d+) (\d+)"', svg)
    return (int(m.group(1)), int(m.group(2)))


def _bar(slide, l, t, w, rate, color=DONE):
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.16))
    bg.fill.solid(); bg.fill.fore_color.rgb = FUT_BG; bg.line.fill.background()
    if rate > 0:
        fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w * rate / 100), Inches(0.16))
        fg.fill.solid(); fg.fill.fore_color.rgb = color; fg.line.fill.background()


def build_deck(members: list[dict], out_path: Path):
    today = dt.date.today()
    iso = today.isocalendar()
    week = f"{iso[0]}-W{iso[1]:02d}"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    from pptx.enum.shapes import MSO_SHAPE

    # ── 表紙＋ダッシュボード ─────────────────
    s = prs.slides.add_slide(blank)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
    band.fill.solid(); band.fill.fore_color.rgb = GOAL; band.line.fill.background()
    _text(s, 0.6, 0.3, 12, 0.7, [("週次ゴール進捗", 30, WHITE, True, PP_ALIGN.LEFT)])
    _text(s, 0.6, 1.0, 12, 0.4, [(f"{week} ／ {today.isoformat()} ／ 対象 {len(members)}名", 13, WHITE, False, PP_ALIGN.LEFT)])

    # ダッシュボード行
    y = 1.95
    _text(s, 0.6, y, 3.0, 0.3, [("メンバー", 12, SUB, True, PP_ALIGN.LEFT)])
    _text(s, 3.6, y, 3.6, 0.3, [("テーマ", 12, SUB, True, PP_ALIGN.LEFT)])
    _text(s, 7.2, y, 2.0, 0.3, [("現在ステージ", 12, SUB, True, PP_ALIGN.LEFT)])
    _text(s, 9.3, y, 3.4, 0.3, [("達成率", 12, SUB, True, PP_ALIGN.LEFT)])
    y += 0.45
    for d in members:
        rate = _rate(d)
        cur = int(d.get("currentStage", 1))
        _text(s, 0.6, y, 3.0, 0.4, [(d.get("name", ""), 14, INK, True, PP_ALIGN.LEFT)])
        _text(s, 3.6, y, 3.6, 0.4, [(d.get("theme", ""), 11, INK, False, PP_ALIGN.LEFT)])
        _text(s, 7.2, y, 2.0, 0.4, [(STAGE_NAMES[cur - 1], 11, NOW, True, PP_ALIGN.LEFT)])
        _text(s, 11.7, y - 0.02, 1.0, 0.4, [(f"{rate}%", 13, DONE, True, PP_ALIGN.LEFT)])
        _bar(s, 9.3, y + 0.12, 2.3, rate)
        y += 0.62

    # ── 各メンバー1枚 ───────────────────────
    for d in members:
        s = prs.slides.add_slide(blank)
        head = d.get("name", "")
        if d.get("note"):
            head += f"（{d['note']}）"
        _text(s, 0.5, 0.3, 12.3, 0.6, [(f"{head}｜{d.get('theme','')}", 22, INK, True, PP_ALIGN.LEFT)])

        # 図解PNG（左）
        png = _png(d)
        w_px, h_px = _png_size(d)
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:
            tf.write(png); tmp = tf.name
        img_h = 6.3
        img_w = img_h * (w_px / h_px)
        if img_w > 8.2:
            img_w = 8.2; img_h = img_w * (h_px / w_px)
        s.shapes.add_picture(tmp, Inches(0.5), Inches(1.05), height=Inches(img_h))

        # サイドパネル（右）
        rx = 8.9
        rate = _rate(d)
        _text(s, rx, 1.1, 3.9, 0.3, [("達成率", 12, SUB, True, PP_ALIGN.LEFT)])
        _text(s, rx, 1.35, 3.9, 0.6, [(f"{rate}%", 34, DONE, True, PP_ALIGN.LEFT)])
        _bar(s, rx, 2.1, 3.9, rate)
        panel = [("ゴール", 12, GOAL, True, PP_ALIGN.LEFT), (d.get("goal", ""), 13, INK, False, PP_ALIGN.LEFT)]
        if d.get("focus"):
            panel += [("今週のフォーカス", 12, NOW, True, PP_ALIGN.LEFT), (d["focus"], 13, INK, False, PP_ALIGN.LEFT)]
        if d.get("next"):
            panel += [("ネクストアクション", 12, SUB, True, PP_ALIGN.LEFT), (d["next"], 13, INK, False, PP_ALIGN.LEFT)]
        _text(s, rx, 2.5, 3.9, 4.4, panel)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path, week


def main(argv):
    args = [a for a in argv[1:] if not a.startswith("-")]
    opts = argv[1:]
    mdir = HERE / "members"
    if args:
        files = [mdir / f"{a}.json" for a in args]
    else:
        files = sorted(p for p in mdir.glob("*.json") if not p.name.startswith("_"))
    members = [load_member_data(str(f)) for f in files if f.exists()]
    if not members:
        print("members/*.json が見つかりません")
        return 1

    out = None
    if "-o" in opts:
        out = Path(opts[opts.index("-o") + 1])
    if out is None:
        iso = dt.date.today().isocalendar()
        out = HERE / "out" / f"週次ゴール進捗_{iso[0]}-W{iso[1]:02d}.pptx"

    path, week = build_deck(members, out)
    print(f"wrote {path}  （{week} / {len(members)}名）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
