#!/usr/bin/env python3
"""人ごとのタスクの繋がり（指示者 → 作業者）をPPTX 1枚に図解する。
左=指示ネットワーク図、右=関係ごとのタスク一覧。データはSupabaseの実データ。
出力: 親フォルダ直下の E.L.P_指示マップ.pptx
"""
import json
import os
import urllib.parse
import urllib.request
from collections import defaultdict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ENV_PATH = os.path.join(os.path.dirname(__file__), "..", ".env")
OUT = os.path.join(os.path.dirname(__file__), "..", "..", "E.L.P_指示マップ.pptx")
JP = "Hiragino Sans"

C_TOP = RGBColor(0x06, 0x4E, 0x3B)
C_LEAD = RGBColor(0x04, 0x7A, 0x55)
C_OP = RGBColor(0x0F, 0x8A, 0x66)
C_ARROW = RGBColor(0x10, 0x99, 0x6B)
C_CHIP = RGBColor(0x10, 0x99, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x06, 0x4E, 0x3B)

# 既知メンバーの配置（中心座標とロール・色）
NODES = {
    "松尾": dict(cx=3.6, cy=1.75, role="代表取締役", color=C_TOP),
    "東":   dict(cx=2.25, cy=3.7, role="車トラブル事業", color=C_LEAD),
    "池田": dict(cx=5.35, cy=3.7, role="清掃事業", color=C_LEAD),
    "鈴木": dict(cx=1.45, cy=5.7, role="オペレーター", color=C_OP),
    "藤田": dict(cx=4.35, cy=5.7, role="オペレーター", color=C_OP),
}
NW, NH = 1.55, 0.8


def load_env():
    env = {}
    for line in open(ENV_PATH, encoding="utf-8"):
        line = line.rstrip("\n")
        if "=" in line and not line.startswith("#"):
            k, v = line.split("=", 1)
            env[k.strip()] = v.strip()
    return env


def fetch_edges():
    env = load_env()
    url, key = env["ELP_SUPABASE_URL"], env["ELP_SUPABASE_SERVICE_ROLE_KEY"]
    H = {"apikey": key, "Authorization": f"Bearer {key}"}
    sel = "title,worker:assignee_id(name),director:director_id(name)"
    q = ("tasks?select=" + urllib.parse.quote(sel)
         + "&director_id=not.is.null&assignee_id=not.is.null&status=neq." + urllib.parse.quote("完了"))
    d = json.load(urllib.request.urlopen(urllib.request.Request(f"{url}/rest/v1/{q}", headers=H), timeout=30))
    edges, selfd = defaultdict(list), defaultdict(list)
    for t in d:
        w = (t.get("worker") or {}).get("name")
        di = (t.get("director") or {}).get("name")
        if not w or not di:
            continue
        (selfd[w] if w == di else edges[(di, w)]).append(t["title"])
    return edges, selfd


def _font(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = JP
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, (s, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(1)
        run = p.add_run(); run.text = s
        _font(run, size, color, bold)
    return tf


def node(slide, name, info, self_count):
    x, y = info["cx"] - NW / 2, info["cy"] - NH / 2
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(NW), Inches(NH))
    sp.fill.solid(); sp.fill.fore_color.rgb = info["color"]; sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.12
    except Exception:
        pass
    label = name + (f"（自己{self_count}）" if self_count else "")
    textbox(slide, x, y + 0.04, NW, NH, [
        (label, 15, WHITE, True),
        (info["role"], 9.5, RGBColor(0xCD, 0xEA, 0xDF), False),
    ], align=PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, w=2.25):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = C_ARROW; c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def chip(slide, mx, my, n):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx - 0.16), Inches(my - 0.16), Inches(0.32), Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = C_CHIP; s.line.color.rgb = WHITE; s.line.width = Pt(1)
    s.shadow.inherit = False
    textbox(slide, mx - 0.16, my - 0.17, 0.32, 0.32, [(str(n), 11, WHITE, True)], align=PP_ALIGN.CENTER)


def edge_anchor(a, b):
    ax, ay = NODES[a]["cx"], NODES[a]["cy"]
    bx, by = NODES[b]["cx"], NODES[b]["cy"]
    if abs(ay - by) < 0.6:
        if ax < bx:
            return ax + NW / 2, ay, bx - NW / 2, by
        return ax - NW / 2, ay, bx + NW / 2, by
    if ay < by:
        return ax, ay + NH / 2, bx, by - NH / 2
    return ax, ay - NH / 2, bx, by + NH / 2


def main():
    edges, selfd = fetch_edges()
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textbox(slide, 0.5, 0.25, 12.3, 0.6, [("人ごとのタスクの繋がり（指示 → 作業）", 24, DARK, True)])
    textbox(slide, 0.5, 0.8, 7.0, 0.35, [("矢印＝指示者から作業者へ ・ 数字＝タスク件数", 12, GRAY, False)])

    for (di, w), tasks in edges.items():
        if di in NODES and w in NODES:
            x1, y1, x2, y2 = edge_anchor(di, w)
            arrow(slide, x1, y1, x2, y2)
            chip(slide, (x1 + x2) / 2, (y1 + y2) / 2, len(tasks))

    for name, info in NODES.items():
        node(slide, name, info, len(selfd.get(name, [])))

    px, pw = 7.9, 5.0
    textbox(slide, px, 1.2, pw, 0.4, [("指示の関係（誰が → 誰に）", 14, DARK, True)])
    y = 1.75
    order = sorted(edges.items(), key=lambda kv: (-len(kv[1]), kv[0]))
    for (di, w), tasks in order:
        textbox(slide, px, y, pw, 0.32, [(f"{di} → {w}（{len(tasks)}件）", 12, C_LEAD, True)])
        y += 0.34
        textbox(slide, px + 0.15, y, pw - 0.15, 0.6, [("・" + " / ".join(tasks), 10.5, GRAY, False)],
                anchor=MSO_ANCHOR.TOP)
        y += 0.52
    selfs = [f"{n}: {' / '.join(ts)}" for n, ts in selfd.items()]
    if selfs:
        textbox(slide, px, y + 0.05, pw, 0.3, [("自己管理タスク", 12, GRAY, True)])
        y += 0.36
        for s in selfs:
            textbox(slide, px + 0.15, y, pw - 0.15, 0.3, [("・" + s, 10.5, GRAY, False)])
            y += 0.3

    prs.save(OUT)
    print(f"saved: {os.path.abspath(OUT)}")


if __name__ == "__main__":
    main()
