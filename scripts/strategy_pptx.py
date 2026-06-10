#!/usr/bin/env python3
"""E.L.P 事業戦略マップ（目的 → 今季の目標 → KPI → タスク）を1枚のPPTXに描く。

部下への説明用。上から下へ「なぜ→何を→どの数字→日々の手」が一目で繋がる。
数値は仮置き（実数に差し替え可）。タスク名は業務ポータルの実データに対応。
出力: 親フォルダ直下の E.L.P_戦略マップ.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = os.path.join(os.path.dirname(__file__), "..", "..", "E.L.P_戦略マップ.pptx")
JP = "Hiragino Sans"

# カラー（エメラルド系・アプリと統一）
C_PURPOSE = RGBColor(0x06, 0x4E, 0x3B)
C_GOAL = RGBColor(0x04, 0x7A, 0x55)
C_KPI = RGBColor(0x0F, 0x8A, 0x66)
C_TASK_FILL = RGBColor(0xEC, 0xFD, 0xF5)
C_TASK_TEXT = RGBColor(0x06, 0x4E, 0x3B)
C_TAG = RGBColor(0x10, 0x99, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)


def _set_font(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = JP
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", JP)


def box(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.08
    except Exception:
        pass
    return sp


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (s, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = s
        _set_font(run, size, color, bold)
    return tb


def tag(slide, x, y, w, label):
    box(slide, x, y, w, 0.42, C_TAG)
    text(slide, x, y, w, 0.42, [(label, 13, WHITE, True)], align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # タイトル
    text(slide, 0.5, 0.25, 12.3, 0.7, [
        ("E.L.P 事業戦略マップ", 26, C_PURPOSE, True),
    ], anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 0.5, 0.78, 12.3, 0.35, [
        ("目的 → 今季の目標 → 追う数字 → 日々のタスク", 13, GRAY, False),
    ])

    LX, LW = 0.5, 1.5            # 左ラベル列
    CX, CW = 2.2, 10.6           # コンテンツ列

    # 1) 目的
    tag(slide, LX, 1.3, LW, "目的")
    box(slide, CX, 1.3, CW, 0.95, C_PURPOSE)
    text(slide, CX + 0.2, 1.3, CW - 0.4, 0.95, [
        ("車のトラブルと建物の維持で、地域の『困った』を最短・確実に解決する", 18, WHITE, True),
    ])

    # 2) 今季の目標
    tag(slide, LX, 2.5, LW, "今季の目標")
    box(slide, CX, 2.5, CW, 0.95, C_GOAL)
    text(slide, CX + 0.2, 2.5, CW - 0.4, 0.95, [
        ("車トラブル受付センターを主力に、月間成約120件・売上1,500万円を達成（2026 Q2）", 16, WHITE, True),
    ])

    # 3) 追う数字（KPI）— 4列
    tag(slide, LX, 3.7, LW, "追う数字")
    kpis = [
        ("リード数（着信）", "月 1,500件"),
        ("成約率", "35% → 40%"),
        ("対応スピード", "受付→手配 5分以内"),
        ("清掃 継続率", "95%以上"),
    ]
    gap = 0.25
    kw = (CW - gap * 3) / 4
    kx = []
    for i, (name, val) in enumerate(kpis):
        x = CX + i * (kw + gap)
        kx.append(x)
        box(slide, x, 3.7, kw, 1.0, C_KPI)
        text(slide, x, 3.74, kw, 1.0, [
            (name, 13, WHITE, True),
            (val, 16, WHITE, True),
        ], align=PP_ALIGN.CENTER)

    # 4) タスク — 各KPIの下に対応する実タスク
    tag(slide, LX, 5.0, LW, "タスク")
    tasks = [
        ["あさひ自動車販売へ協定書送付", "レッカー業者の契約更新"],
        ["夜間対応マニュアル整備", "受付フローのKPI集計"],
        ["顧客管理システム刷新（ベンダー選定）", "既存データ移行計画"],
        ["グリーンビル契約更新交渉", "清掃スタッフのシフト表作成"],
    ]
    for i, items in enumerate(tasks):
        x = kx[i]
        box(slide, x, 5.0, kw, 1.7, C_TASK_FILL)
        lines = [(f"• {it}", 11.5, C_TASK_TEXT, False) for it in items]
        text(slide, x + 0.05, 5.05, kw - 0.1, 1.6, lines, anchor=MSO_ANCHOR.TOP)

    # 縦の流れを示す下向き矢印（左ラベル列の中央）
    ax = LX + LW / 2 - 0.12
    for y0 in (2.28, 3.48, 4.78):
        box(slide, ax, y0, 0.24, 0.2, C_TAG, shape=MSO_SHAPE.DOWN_ARROW)

    # 注記
    text(slide, CX, 6.85, CW, 0.35, [
        ("※ 目標・KPIの数値は仮置きです。実数に差し替えてください。タスクは業務ポータルの実データに対応。", 10.5, GRAY, False),
    ])

    prs.save(OUT)
    print(f"saved: {os.path.abspath(OUT)}")


if __name__ == "__main__":
    main()
