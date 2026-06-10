#!/usr/bin/env python3
"""E.L.P 事業戦略 全体図（ダミー）を複数ページのPPTXで作る。
ホワイトボード（売上20億→70億、会長→社長→部長→社員、提携/訪問/自社レンタカー）を反映。
構成: 表紙 / 全体像 / 体制図 / 事業部詳細×3
出力: 親フォルダ直下の E.L.P_事業戦略_全体図.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = os.path.join(os.path.dirname(__file__), "..", "..", "E.L.P_事業戦略_全体図.pptx")
JP = "Hiragino Sans"

DARK = RGBColor(0x06, 0x4E, 0x3B)
GOAL = RGBColor(0x04, 0x7A, 0x55)
KPI = RGBColor(0x0F, 0x8A, 0x66)
LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
TAG = RGBColor(0x10, 0x99, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LGREEN = RGBColor(0xCD, 0xEA, 0xDF)
BLUEBG = RGBColor(0xE6, 0xF1, 0xFB)
BLUETX = RGBColor(0x0C, 0x44, 0x7C)
LEVEL = {"大": (RGBColor(0xD1, 0xFA, 0xE5), DARK), "中": (BLUEBG, BLUETX), "小": (RGBColor(0xF1, 0xEF, 0xE8), GRAY)}

W, H = 13.333, 7.5


def _font(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = JP
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP)


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (s, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(2)
        run = p.add_run(); run.text = s
        _font(run, size, color, bold)
    return tf


def box(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.08
    except Exception:
        pass
    return sp


def arrow(slide, x1, y1, x2, y2, color=TAG, w=2.25):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(s, title, sub=None):
    box(s, 0, 0, W, 0.12, TAG, shape=MSO_SHAPE.RECTANGLE)
    text(s, 0.5, 0.28, 11.5, 0.6, [(title, 24, DARK, True)])
    if sub:
        text(s, 0.5, 0.85, 11.5, 0.35, [(sub, 12.5, GRAY, False)])


def s_cover(prs):
    s = slide(prs)
    box(s, 0, 0, W, H, DARK, shape=MSO_SHAPE.RECTANGLE)
    text(s, 0.9, 1.5, 11.5, 0.5, [("E.L.P 事業戦略 全体図", 16, LGREEN, False)])
    text(s, 0.9, 2.1, 11.5, 1.4, [("売上 20億円", 60, WHITE, True)])
    text(s, 0.9, 3.5, 11.5, 0.6, [("— その先に見る景色は、70億円 —", 20, LGREEN, False)])
    text(s, 0.9, 4.6, 11.5, 0.5, [("2026年度 ・ イノベラ連携を軸に主力事業を伸ばす", 14, LGREEN, False)])
    text(s, 0.9, 5.4, 11.5, 0.5, [("会長 → 社長 → 部長 → 社員 が一本で繋がる実行体制", 14, WHITE, False)])
    text(s, 0.9, 6.7, 11.5, 0.4, [("※ 本資料の数値・体制はダミー（叩き台）です。", 11, LGREEN, False)])


def s_overview(prs):
    s = slide(prs)
    header(s, "全体像：KGI → 事業部 → KPI", "売上20億を4つの事業部で分担し、各事業部のKPIに落とす")
    box(s, 4.0, 1.35, 5.3, 0.85, DARK)
    text(s, 4.0, 1.35, 5.3, 0.85, [("KGI：売上20億円（→70億）", 18, WHITE, True)], align=PP_ALIGN.CENTER)
    divs = [
        ("車トラブル受付センター事業", ["成約率 35→40%", "リード 月1,500件"]),
        ("提携・新規開拓", ["提携工場 月10社", "訪問 100件/月"]),
        ("清掃事業", ["契約継続率 95%"]),
        ("自社レンタカー事業", ["稼働率 70%", "提携先 5社"]),
    ]
    n = len(divs); gap = 0.3; bw = (12.3 - gap * (n - 1)) / n; x0 = 0.5
    for i, (name, kpis) in enumerate(divs):
        x = x0 + i * (bw + gap)
        arrow(s, 6.65, 2.2, x + bw / 2, 3.0)
        box(s, x, 3.0, bw, 0.7, GOAL)
        text(s, x, 3.0, bw, 0.7, [(name, 12.5, WHITE, True)], align=PP_ALIGN.CENTER)
        box(s, x, 3.85, bw, 1.9, LIGHT)
        lines = [("追う数字 (KPI)", 10, GRAY, True)] + [("・" + k, 12, DARK, False) for k in kpis]
        text(s, x + 0.1, 3.95, bw - 0.2, 1.7, lines, anchor=MSO_ANCHOR.TOP)
    text(s, 0.5, 6.9, 12.3, 0.4, [("各事業部の詳細（大・中・小タスクと担当）は次ページ以降。", 11, GRAY, False)])


def s_org(prs):
    s = slide(prs)
    header(s, "実行体制：会長 → 社長 → 部長 → 社員", "上位の目標が、現場の一人ひとりの動きまで降りてくる")
    tiers = [
        ("会長", "経営の最終意思決定", ["全社の方向づけ", "投資判断"], DARK),
        ("社長（松尾）", "全社の数字に責任", ["個人3 / 法人4 を月次で", "事業部KPIを統括"], GOAL),
        ("部長（東・池田）", "事業部のKPI達成", ["社員へリスト配布（期限明確）", "訪問100件/月を設計"], KPI),
        ("社員（鈴木・藤田 他）", "実行・現場", ["テレアポ・訪問", "受付/清掃/レンタカー現場"], RGBColor(0x3B, 0x9E, 0x7E)),
    ]
    n = len(tiers); gap = 0.4; bw = (12.3 - gap * (n - 1)) / n; x0 = 0.5; y = 1.7
    for i, (role, resp, items, color) in enumerate(tiers):
        x = x0 + i * (bw + gap)
        if i > 0:
            arrow(s, x - gap - 0.02, y + 1.6, x + 0.02, y + 1.6)
        box(s, x, y, bw, 0.95, color)
        text(s, x, y + 0.05, bw, 0.95, [(role, 15, WHITE, True), (resp, 10, LGREEN, False)], align=PP_ALIGN.CENTER)
        box(s, x, y + 1.1, bw, 2.4, LIGHT)
        text(s, x + 0.12, y + 1.22, bw - 0.24, 2.2,
             [("役割", 10, GRAY, True)] + [("・" + it, 12, DARK, False) for it in items], anchor=MSO_ANCHOR.TOP)
    box(s, 0.5, 5.5, 12.3, 1.2, BLUEBG)
    text(s, 0.7, 5.6, 12.0, 1.0, [
        ("重点テーマ：自社レンタカー業務推進", 13, BLUETX, True),
        ("レッカー対応とレンタカーをセット提供し、客単価と継続率を引き上げる新規の柱。", 12, BLUETX, False),
    ], anchor=MSO_ANCHOR.TOP)


def s_division(prs, title, lead, kpis, tasks):
    s = slide(prs)
    header(s, f"事業部詳細：{title}", f"事業部長：{lead}")
    n = len(kpis); gap = 0.3; bw = (12.3 - gap * (n - 1)) / n; x0 = 0.5; ky = 1.4
    for i, (name, cur, tgt) in enumerate(kpis):
        x = x0 + i * (bw + gap)
        box(s, x, ky, bw, 0.95, KPI)
        text(s, x, ky + 0.05, bw, 0.95, [(name, 12.5, WHITE, True), (f"{cur} → {tgt}", 13, WHITE, True)],
             align=PP_ALIGN.CENTER)
    text(s, 0.5, 2.6, 12.3, 0.35, [("タスク（大 → 中 → 小） ・ 担当＝作業者 / 指示者", 12, GRAY, True)])
    y = 3.05
    indent = {"大": 0.5, "中": 1.3, "小": 2.1}
    for level, ttl, worker, director in tasks:
        x = indent[level]
        fill, tx = LEVEL[level]
        box(s, x, y, 0.55, 0.4, fill)
        text(s, x, y, 0.55, 0.4, [(level, 12, tx, True)], align=PP_ALIGN.CENTER)
        text(s, x + 0.7, y, 7.2, 0.4, [(ttl, 13, DARK, True)])
        text(s, 8.3, y, 4.5, 0.4, [(f"作業 {worker}  /  指示 {director}", 11, GRAY, False)], align=PP_ALIGN.RIGHT)
        y += 0.52


def build():
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    s_cover(prs)
    s_overview(prs)
    s_org(prs)
    s_division(
        prs, "車トラブル受付センター事業", "東",
        [("成約率", "35%", "40%"), ("リード数（着信）", "1,200件", "1,500件"), ("平均対応", "8分", "5分")],
        [("大", "受付フローの改善", "東", "松尾"),
         ("中", "夜間対応マニュアル整備", "鈴木", "東"),
         ("小", "深夜帯の一次対応スクリプト作成", "鈴木", "東"),
         ("中", "受付フローのKPI集計", "藤田", "東"),
         ("大", "提携先の拡大", "東", "松尾"),
         ("中", "あさひ自動車販売へ協定書送付", "東", "東")],
    )
    s_division(
        prs, "提携・新規開拓（工場 / テレアポ / 訪問）", "東",
        [("提携工場", "6社", "月10社"), ("訪問", "60件", "100件/月"), ("テレアポ", "600件", "1,000件/月")],
        [("大", "新規提携の獲得", "東", "松尾"),
         ("中", "工場リストの作成・社員へ配布（期限明確）", "藤田", "東"),
         ("小", "Aさんへリストを渡す（◯月◯日まで）", "藤田", "東"),
         ("中", "テレアポ → 訪問アポ獲得", "鈴木", "東"),
         ("小", "訪問100件にかける", "鈴木", "東")],
    )
    s_division(
        prs, "清掃事業 ＆ 自社レンタカー事業", "池田",
        [("清掃 継続率", "92%", "95%"), ("レンタカー稼働率", "—", "70%"), ("レンタカー提携", "1社", "5社")],
        [("大", "既存顧客の維持・拡大（清掃）", "池田", "松尾"),
         ("中", "グリーンビル契約更新交渉", "池田", "池田"),
         ("中", "清掃スタッフのシフト表作成", "藤田", "池田"),
         ("大", "自社レンタカー業務の立ち上げ", "池田", "松尾"),
         ("中", "レンタカー運用フロー構築", "鈴木", "池田"),
         ("中", "レッカー＋レンタカーのセット提案", "藤田", "池田")],
    )
    prs.save(OUT)
    print(f"saved: {os.path.abspath(OUT)} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
